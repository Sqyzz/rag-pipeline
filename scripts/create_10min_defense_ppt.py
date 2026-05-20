from __future__ import annotations

import io
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "论文改稿" / "华中师范大学硕博PPT模板.pptx"
OUT = ROOT / "论文改稿" / "沈秋雨-基于GraphRAG的企业知识库检索优化研究-10分钟汇报.pptx"


COLORS = {
    "teal": RGBColor(0x00, 0x63, 0x69),
    "deep_teal": RGBColor(0x00, 0x46, 0x50),
    "light_teal": RGBColor(0xE5, 0xF2, 0xF1),
    "gold": RGBColor(0xC2, 0x97, 0x48),
    "red": RGBColor(0x9B, 0x1B, 0x30),
    "ink": RGBColor(0x19, 0x24, 0x2D),
    "muted": RGBColor(0x5C, 0x6B, 0x73),
    "line": RGBColor(0xD9, 0xE2, 0xE1),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "bg": RGBColor(0xF8, 0xFA, 0xFA),
}


FONT_CN = "Microsoft YaHei"
FONT_SERIF = "SimSun"


def delete_all_slides(prs: Presentation) -> None:
    sld_id_lst = prs.slides._sldIdLst  # noqa: SLF001 - python-pptx has no public delete API.
    for sld_id in list(sld_id_lst):
        prs.part.drop_rel(sld_id.rId)
        sld_id_lst.remove(sld_id)


def set_bg(slide, color="bg") -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS[color]


def add_box(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    font_size: int = 18,
    color: str = "ink",
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    fill: str | None = None,
    line: str | None = None,
    radius=MSO_SHAPE.RECTANGLE,
    margin: float = 0.12,
):
    shape = slide.shapes.add_shape(radius, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS[fill]
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = COLORS[line]
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin / 1.5)
    tf.margin_bottom = Inches(margin / 1.5)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT_CN
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = COLORS[color]
    return shape


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    font_size: int = 18,
    color: str = "ink",
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    font=FONT_CN,
):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = COLORS[color]
    return tb


def add_bullets(
    slide,
    items: list[str],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    font_size: int = 15,
    color: str = "ink",
    bullet_color: str = "gold",
    gap_after: int = 8,
):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(gap_after)
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        r1 = p.add_run()
        r1.text = "• "
        r1.font.name = FONT_CN
        r1.font.size = Pt(font_size)
        r1.font.bold = True
        r1.font.color.rgb = COLORS[bullet_color]
        r2 = p.add_run()
        r2.text = item
        r2.font.name = FONT_CN
        r2.font.size = Pt(font_size)
        r2.font.color.rgb = COLORS[color]
    return tb


def add_line(slide, x1, y1, x2, y2, color="line", width=1.2):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = COLORS[color]
    line.line.width = Pt(width)
    return line


def add_logo(slide, prs: Presentation, x=0.55, y=0.28, w=1.95) -> None:
    with zipfile.ZipFile(TEMPLATE) as zf:
        data = zf.read("ppt/media/image3.png")
    slide.shapes.add_picture(io.BytesIO(data), Inches(x), Inches(y), width=Inches(w))


def add_header(slide, prs: Presentation, title: str, page: int, total: int = 12) -> None:
    add_logo(slide, prs, 0.42, 0.18, 1.58)
    add_text(slide, title, 2.15, 0.24, 8.4, 0.42, font_size=17, color="deep_teal", bold=True)
    add_text(slide, f"{page:02d}/{total:02d}", 11.55, 0.24, 1.0, 0.32, font_size=10, color="muted", align=PP_ALIGN.RIGHT)
    add_line(slide, 0.42, 0.78, 12.95, 0.78, color="line", width=1)


def add_footer(slide) -> None:
    add_text(slide, "基于GraphRAG的企业知识库检索优化研究", 0.55, 7.12, 5.5, 0.22, font_size=8, color="muted")


def title_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide, "bg")
    add_logo(slide, prs, 0.72, 0.45, 2.25)
    add_box(slide, "", 0, 0, 13.33, 0.11, fill="teal")
    add_box(slide, "", 0.0, 6.95, 13.33, 0.55, fill="deep_teal")
    add_text(slide, "基于GraphRAG的企业知识库检索优化研究", 0.9, 1.7, 10.6, 0.78, font_size=28, color="deep_teal", bold=True, font=FONT_SERIF)
    add_text(slide, "10分钟论文汇报", 0.95, 2.58, 4.2, 0.42, font_size=18, color="gold", bold=True)
    add_line(slide, 0.95, 3.18, 6.1, 3.18, color="gold", width=2.2)
    add_text(slide, "汇报人：沈秋雨", 0.95, 3.48, 4.7, 0.28, font_size=14, color="ink")
    add_text(slide, "指导教师：程秀峰 教授", 0.95, 3.82, 4.7, 0.28, font_size=14, color="ink")
    add_text(slide, "学科专业：图书情报", 0.95, 4.16, 4.7, 0.28, font_size=14, color="ink")
    add_text(slide, "研究方向：信息组织与检索", 0.95, 4.50, 4.7, 0.28, font_size=14, color="ink")
    add_text(slide, "华中师范大学信息管理学院\n2026年5月", 0.95, 5.25, 4.2, 0.62, font_size=13, color="muted")

    # Decorative knowledge graph.
    cx, cy = 9.55, 3.55
    nodes = [
        (cx - 1.6, cy - 0.9, "事实"),
        (cx, cy - 1.35, "实体"),
        (cx + 1.65, cy - 0.55, "关系"),
        (cx - 1.05, cy + 0.7, "主题"),
        (cx + 1.25, cy + 0.95, "证据"),
    ]
    for a, b in [(0, 1), (1, 2), (1, 3), (2, 4), (3, 4), (0, 3)]:
        add_line(slide, nodes[a][0] + 0.35, nodes[a][1] + 0.18, nodes[b][0] + 0.35, nodes[b][1] + 0.18, color="line", width=1.8)
    for x, y, label in nodes:
        add_box(slide, label, x, y, 0.82, 0.42, font_size=12, color="white", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="teal", radius=MSO_SHAPE.ROUNDED_RECTANGLE, margin=0.02)
    return slide


def contents_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide, "bg")
    add_header(slide, prs, "目录", 2)
    add_text(slide, "CONTENTS", 0.78, 1.16, 2.8, 0.36, font_size=12, color="gold", bold=True)
    add_text(slide, "汇报内容", 0.75, 1.48, 2.4, 0.46, font_size=22, color="deep_teal", bold=True)
    sections = [
        ("01", "研究背景", "企业知识检索为什么难"),
        ("02", "核心思路", "从企业知识架构出发重构GraphRAG"),
        ("03", "方法设计", "四层图结构与Agentic Decomposer联动"),
        ("04", "实验验证", "数据集构建与关键结果"),
        ("05", "结论展望", "适用场景与后续方向"),
    ]
    y = 2.0
    for no, title, desc in sections:
        add_text(slide, no, 1.0, y + 0.04, 0.7, 0.34, font_size=17, color="gold", bold=True)
        add_text(slide, title, 1.85, y, 1.8, 0.36, font_size=16, color="deep_teal", bold=True)
        add_text(slide, desc, 3.35, y + 0.03, 6.8, 0.34, font_size=13, color="ink")
        add_line(slide, 1.0, y + 0.56, 10.8, y + 0.56, color="line")
        y += 0.85
    add_box(slide, "主线：企业知识天然分层 → 图谱分层组织 → Query规划路由 → 多粒度证据召回", 1.0, 6.35, 10.9, 0.5, font_size=14, color="white", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="deep_teal", radius=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_footer(slide)
    return slide


def background_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide, "bg")
    add_header(slide, prs, "研究背景：企业知识检索的核心矛盾", 3)
    add_text(slide, "企业知识不同于普通网页文本，检索系统需要同时处理三个粒度。", 0.72, 1.1, 9.5, 0.36, font_size=15, color="muted")
    cards = [
        ("细粒度事实密集", "金额、期限、义务、责任主体\n需要条款级精准定位", "Local"),
        ("跨条款关系复杂", "触发条件、例外条款、责任链\n需要显式结构推理", "Structural"),
        ("全局语义需整合", "合同风险、义务分布、主题画像\n需要宏观归纳能力", "Global"),
    ]
    x = 0.85
    for title, body, tag in cards:
        add_box(slide, tag, x, 1.82, 1.35, 0.34, font_size=11, color="white", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="gold", radius=MSO_SHAPE.ROUNDED_RECTANGLE, margin=0.02)
        add_box(slide, title, x, 2.28, 3.45, 0.55, font_size=17, color="deep_teal", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="white", line="line", radius=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_box(slide, body, x, 2.95, 3.45, 1.1, font_size=13, color="ink", align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="white", line="line", radius=MSO_SHAPE.ROUNDED_RECTANGLE)
        x += 4.0
    add_box(slide, "传统VectorRAG：能找相似文本，但缺少结构关系。\n原生GraphRAG：能做全局社区摘要，但容易弱化底层事实定位。", 1.0, 4.72, 5.35, 1.08, font_size=14, color="ink", fill="light_teal", line="line", radius=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_box(slide, "本文问题：如何在企业文档中同时保持局部证据精确性、结构关系可推理性和全局语义整合能力？", 6.78, 4.72, 5.45, 1.08, font_size=15, color="white", bold=True, fill="deep_teal", radius=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_footer(slide)
    return slide


def overall_idea_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide, "bg")
    add_header(slide, prs, "研究目标：从企业知识架构出发重构GraphRAG", 4)
    add_text(slide, "E-GraphRAG不是简单“加图谱”，而是将知识组织和检索调度同时改造。", 0.75, 1.1, 10.2, 0.35, font_size=15, color="muted")
    steps = [
        ("企业文档", "合同、制度、SOP"),
        ("四层知识图谱", "事实-实体-关系-主题"),
        ("Query规划", "意图拆解与需求编译"),
        ("三类路由", "Local / Structural / Global"),
        ("证据回答", "可追溯、低噪声"),
    ]
    x = 0.78
    y = 2.25
    for idx, (title, desc) in enumerate(steps):
        add_box(slide, title, x, y, 1.95, 0.55, font_size=15, color="white", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="teal", radius=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_box(slide, desc, x, y + 0.7, 1.95, 0.65, font_size=11, color="ink", align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="white", line="line", radius=MSO_SHAPE.ROUNDED_RECTANGLE)
        if idx < len(steps) - 1:
            add_line(slide, x + 1.98, y + 0.28, x + 2.42, y + 0.28, color="gold", width=2)
        x += 2.45
    add_box(slide, "底层：四层结构解决“企业知识如何被组织”\n顶层：Agentic Decomposer解决“用户问题应进入哪一层检索”", 1.15, 4.62, 10.9, 0.95, font_size=17, color="deep_teal", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="white", line="gold", radius=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_bullets(slide, ["目标不是最大化召回，而是在低容错企业场景中提升证据纯度、逻辑稳定性和可追溯性。"], 1.35, 6.0, 10.1, 0.42, font_size=14)
    add_footer(slide)
    return slide


def four_layers_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide, "bg")
    add_header(slide, prs, "为什么要分四层结构", 5)
    layers = [
        ("社区层", "宏观主题：风险类型、义务群组、争议解决机制", "Global"),
        ("关键词/关系层", "中观连接：条款关系、主题词、跨片段路径", "Structural"),
        ("实体层", "业务对象：合同主体、条款对象、责任主体", "Local / Structural"),
        ("属性层", "原子事实：金额、期限、义务、限制条件", "Local"),
    ]
    y = 1.45
    widths = [7.2, 6.5, 5.8, 5.1]
    for i, (name, desc, route) in enumerate(layers):
        w = widths[i]
        x = 0.85 + (7.2 - w) / 2
        color = ["deep_teal", "teal", "light_teal", "white"][i]
        text_color = "white" if i < 2 else "ink"
        add_box(slide, name, x, y, 1.35, 0.52, font_size=15, color=text_color, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill=color, line="line", radius=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_box(slide, desc, x + 1.48, y, w - 2.55, 0.52, font_size=12, color=text_color, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill=color, line="line", radius=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_box(slide, route, x + w - 0.95, y + 0.08, 0.85, 0.34, font_size=8, color="white", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="gold", radius=MSO_SHAPE.ROUNDED_RECTANGLE, margin=0.01)
        y += 0.82
    add_box(slide, "Chunk 原文片段不直接参与拓扑路径计算，而是作为每个节点/关系的底层证据来源，保证可追溯。", 0.98, 4.92, 6.7, 0.58, font_size=13, color="deep_teal", fill="white", line="line", radius=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(slide, "设计理由", 8.45, 1.42, 2.2, 0.36, font_size=18, color="deep_teal", bold=True)
    add_bullets(slide, [
        "企业知识天然具有“微观事实-业务实体-结构关系-宏观主题”的层级。",
        "只保留文本块：能召回但难以推理。",
        "只依赖社区摘要：能概括但容易丢失细节。",
        "四层结构让不同类型的问题进入匹配的证据空间。",
    ], 8.35, 1.95, 4.2, 3.05, font_size=13)
    add_box(slide, "核心：用结构分层把企业知识从碎片文本变成可路由、可推理、可溯源的知识空间。", 8.25, 5.38, 4.25, 0.74, font_size=13, color="white", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="deep_teal", radius=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_footer(slide)
    return slide


def decomposer_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide, "bg")
    add_header(slide, prs, "四层结构如何与Agentic Decomposer联动", 6)
    add_box(slide, "用户Query", 0.7, 1.25, 1.65, 0.55, font_size=15, color="white", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="deep_teal", radius=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_box(slide, "这份合同中有哪些可能影响后续履约风险的条款？", 0.7, 1.95, 3.0, 0.8, font_size=12, color="ink", fill="white", line="line", radius=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_line(slide, 3.78, 2.35, 4.45, 2.35, color="gold", width=2)
    add_box(slide, "Agentic\nDecomposer", 4.48, 1.7, 1.75, 1.25, font_size=15, color="white", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="teal", radius=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_line(slide, 6.28, 2.35, 6.92, 2.35, color="gold", width=2)
    add_box(slide, "结构化检索计划", 6.95, 1.25, 2.2, 0.55, font_size=14, color="white", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="deep_teal", radius=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_bullets(slide, ["route_type", "semantic_need", "anchors", "target_patterns", "检索预算"], 6.95, 1.95, 2.35, 1.65, font_size=12, gap_after=3)
    add_line(slide, 9.35, 2.35, 9.9, 2.35, color="gold", width=2)
    routes = [
        ("Local", "查违约责任、终止条件等具体条款", "属性层 / 实体层"),
        ("Structural", "识别条款之间的条件依赖", "关系层"),
        ("Global", "形成整体风险画像", "社区层"),
    ]
    y = 1.22
    for route, desc, layer in routes:
        add_box(slide, route, 9.95, y, 1.05, 0.42, font_size=11, color="white", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="gold", radius=MSO_SHAPE.ROUNDED_RECTANGLE, margin=0.02)
        add_box(slide, desc + "\n→ " + layer, 11.05, y, 1.65, 0.72, font_size=9, color="ink", align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="white", line="line", radius=MSO_SHAPE.ROUNDED_RECTANGLE, margin=0.04)
        y += 1.0
    add_box(slide, "一句话：四层图结构提供“可检索的知识空间”，Agentic Decomposer决定“问题应该进入哪一层、以什么顺序检索”。", 1.02, 5.35, 11.25, 0.76, font_size=16, color="white", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="deep_teal", radius=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_footer(slide)
    return slide


def rationale_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide, "bg")
    add_header(slide, prs, "为什么这样设计：复合型企业问题需要动态组合路径", 7)
    add_text(slide, "“这份合同有没有风险？”表面是全局问题，实际需要三步证据链。", 0.72, 1.1, 10.5, 0.34, font_size=15, color="muted")
    steps = [
        ("1", "先定位底层条款", "违约责任、保密义务、终止条件、责任限制", "Local"),
        ("2", "再识别条款关系", "触发条件、例外范围、责任链条", "Structural"),
        ("3", "最后形成风险判断", "争议解决、赔偿责任、信息安全风险画像", "Global"),
    ]
    x = 0.85
    for no, title, desc, route in steps:
        add_box(slide, no, x, 2.0, 0.5, 0.5, font_size=17, color="white", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="gold", radius=MSO_SHAPE.OVAL, margin=0.01)
        add_box(slide, title, x + 0.65, 1.88, 2.55, 0.55, font_size=15, color="deep_teal", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="white", line="line", radius=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_box(slide, desc, x, 2.65, 3.2, 1.05, font_size=12, color="ink", align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="white", line="line", radius=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_box(slide, route, x + 1.0, 3.88, 1.2, 0.38, font_size=10, color="white", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="teal", radius=MSO_SHAPE.ROUNDED_RECTANGLE, margin=0.02)
        x += 4.1
    add_box(slide, "直接走全局摘要：可能遗漏关键条款\n只走局部检索：难以形成风险判断\n规划驱动检索：按问题意图动态组合 Local / Structural / Global", 1.08, 5.12, 11.1, 0.92, font_size=15, color="white", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="deep_teal", radius=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_footer(slide)
    return slide


def mapping_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide, "bg")
    add_header(slide, prs, "如何确保用户Query需求被正确映射", 8)
    safeguards = [
        ("任务类型约束", "先归入 Local / Structural / Global，而不是直接全文检索。"),
        ("Schema约束", "结合合同实体类型、关系类型和条款类型生成检索需求。"),
        ("路由字段约束", "每个子问题必须包含 route_type、anchors、target_patterns 等字段。"),
    ]
    x = 0.8
    for i, (title, desc) in enumerate(safeguards, 1):
        add_box(slide, f"{i}", x, 1.65, 0.55, 0.55, font_size=16, color="white", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="gold", radius=MSO_SHAPE.OVAL, margin=0.01)
        add_box(slide, title, x + 0.72, 1.58, 2.45, 0.46, font_size=15, color="deep_teal", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="white", line="line", radius=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_box(slide, desc, x, 2.3, 3.25, 1.05, font_size=12, color="ink", align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="white", line="line", radius=MSO_SHAPE.ROUNDED_RECTANGLE)
        if i < 3:
            add_line(slide, x + 3.35, 2.1, x + 3.78, 2.1, color="gold", width=2)
        x += 4.0
    add_box(slide, "后端编译", 1.1, 4.45, 1.65, 0.5, font_size=14, color="white", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="teal", radius=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_line(slide, 2.82, 4.7, 3.42, 4.7, color="gold", width=2)
    add_box(slide, "结构化检索需求 → 可执行检索动作\nTop-k、搜索半径、目标文档范围、回退机制", 3.48, 4.25, 5.6, 0.92, font_size=14, color="ink", align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="white", line="line", radius=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_line(slide, 9.15, 4.7, 9.75, 4.7, color="gold", width=2)
    add_box(slide, "降低自然语言意图与检索层脱节的风险", 9.82, 4.45, 2.45, 0.5, font_size=12, color="white", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="deep_teal", radius=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_box(slide, "核心保障：任务分类 + Schema约束 + 结构化路由字段", 1.1, 5.86, 11.1, 0.52, font_size=16, color="deep_teal", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="light_teal", line="line", radius=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_footer(slide)
    return slide


def dataset_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide, "bg")
    add_header(slide, prs, "数据集构建：如何确保符合核心研究场景", 9)
    add_text(slide, "本文选择 CUAD 合同数据集作为典型企业文档语料。", 0.72, 1.08, 8.0, 0.34, font_size=15, color="muted")
    flow = [
        ("CUAD合同", "企业级长文档"),
        ("条款切片", "保留Chunk溯源"),
        ("Schema抽取", "实体/属性/关系"),
        ("分层图谱", "四层知识空间"),
        ("测试任务", "Local/Structural/Global"),
    ]
    x = 0.6
    y = 2.02
    for idx, (title, desc) in enumerate(flow):
        add_box(slide, title, x, y, 1.72, 0.52, font_size=13, color="white", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="teal", radius=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_box(slide, desc, x, y + 0.63, 1.72, 0.55, font_size=10, color="ink", align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="white", line="line", radius=MSO_SHAPE.ROUNDED_RECTANGLE)
        if idx < len(flow) - 1:
            add_line(slide, x + 1.77, y + 0.27, x + 2.18, y + 0.27, color="gold", width=2)
        x += 2.28
    add_text(slide, "为什么合同数据符合企业知识检索核心场景", 0.85, 4.15, 6.2, 0.36, font_size=17, color="deep_teal", bold=True)
    add_bullets(slide, [
        "合同是典型企业知识载体，具有明确条款结构和审计需求。",
        "条款内包含金额、期限、主体、义务等细粒度事实。",
        "条款间存在条件、例外、责任边界等结构关系。",
        "合同风险审查需要从多个条款上升到全局语义判断。",
    ], 0.9, 4.65, 6.4, 1.72, font_size=12, gap_after=5)
    add_box(slide, "测试任务围绕三类能力构建：\n单条款业务问题\n条款解释型问题\n多条款联动问题\n合同风险审查", 8.1, 4.25, 3.8, 1.72, font_size=13, color="white", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="deep_teal", radius=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_footer(slide)
    return slide


def results_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide, "bg")
    add_header(slide, prs, "实验结果：按三类企业知识任务观察", 10)
    add_text(slide, "表5-6关键指标摘录；括号内为主要对照基线，用于说明改进来自哪类能力。", 0.75, 1.05, 10.8, 0.3, font_size=12, color="muted")

    cards = [
        (
            "Local 局部事实检索",
            "查具体条款 / 金额 / 期限",
            [
                "Context Precision：0.8765",
                "vs GraphRAG 0.6763 / VectorRAG 0.6480",
                "Answer Correctness：0.7216（最高）",
                "Semantic Similarity：0.9105（最高）",
                "背景噪声：35.6% → 12.3%",
            ],
            "说明：局部路由把问题压到属性/实体证据空间，减少无关条款干扰。",
        ),
        (
            "Structural 结构关系推理",
            "查跨条款关系 / 责任链",
            [
                "Context Precision：0.5250（最高）",
                "vs LightRAG 0.4310 / GraphRAG 0.3745",
                "Semantic Similarity：0.8453（最高）",
                "Context Recall：0.4817（有收敛）",
                "Answer Correctness：0.5643",
            ],
            "说明：路径约束提高证据纯度，但会牺牲部分宽泛召回。",
        ),
        (
            "Global 全局语义整合",
            "查整体风险 / 主题画像",
            [
                "Context Recall：0.7326（最高）",
                "Context Precision：0.6733（最高）",
                "Context MRR：0.9100（最高）",
                "Answer Correctness：0.6232（最高）",
                "Semantic Similarity：0.9381（最高）",
            ],
            "说明：社区层先聚合主题，再下钻证据，提升全局判断的事实边界。",
        ),
    ]
    x = 0.55
    for title, subtitle, rows, explain in cards:
        add_box(slide, title, x, 1.55, 3.9, 0.44, font_size=13, color="white", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="deep_teal", radius=MSO_SHAPE.ROUNDED_RECTANGLE, margin=0.02)
        add_box(slide, subtitle, x, 2.05, 3.9, 0.34, font_size=10, color="deep_teal", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="light_teal", line="line", radius=MSO_SHAPE.ROUNDED_RECTANGLE, margin=0.02)
        add_bullets(slide, rows, x + 0.12, 2.58, 3.65, 1.95, font_size=9, gap_after=2)
        add_box(slide, explain, x, 4.72, 3.9, 0.76, font_size=10, color="ink", align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="white", line="line", radius=MSO_SHAPE.ROUNDED_RECTANGLE, margin=0.05)
        x += 4.15

    add_box(slide, "更谨慎的结论：E-GraphRAG不是所有指标都绝对占优；它的优势主要体现在企业场景更看重的证据精确度、语义一致性、事实边界和可追溯性上。", 0.88, 6.08, 11.55, 0.58, font_size=13, color="white", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="teal", radius=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_footer(slide)
    return slide


def conclusion_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide, "bg")
    add_header(slide, prs, "研究结论与适用场景", 11)
    conclusions = [
        ("结论一", "企业知识检索的关键不是单纯提升召回，而是让不同粒度的问题进入匹配的知识层级。"),
        ("结论二", "四层图结构解决知识组织问题，使微观事实、结构关系和宏观主题能够统一表达。"),
        ("结论三", "Agentic Decomposer解决Query到检索路径的映射问题，动态调度Local、Structural和Global。"),
    ]
    y = 1.28
    for tag, text in conclusions:
        add_box(slide, tag, 0.85, y, 1.2, 0.46, font_size=12, color="white", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="gold", radius=MSO_SHAPE.ROUNDED_RECTANGLE, margin=0.02)
        add_box(slide, text, 2.2, y - 0.05, 9.45, 0.58, font_size=14, color="ink", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE, fill="white", line="line", radius=MSO_SHAPE.ROUNDED_RECTANGLE)
        y += 0.9
    add_text(slide, "适用场景", 0.85, 4.45, 1.8, 0.35, font_size=17, color="deep_teal", bold=True)
    scen = ["合同审查", "制度问答", "企业合规分析", "长文档风险评估"]
    x = 0.9
    for s in scen:
        add_box(slide, s, x, 4.95, 2.0, 0.48, font_size=13, color="white", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="teal", radius=MSO_SHAPE.ROUNDED_RECTANGLE)
        x += 2.3
    add_box(slide, "后续方向：外部行动代理、知识图谱动态演进、多智能体协同评估。", 1.0, 6.05, 10.9, 0.52, font_size=15, color="deep_teal", bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, fill="light_teal", line="line", radius=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_footer(slide)
    return slide


def ending_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide, "deep_teal")
    add_box(slide, "", 0.62, 0.42, 2.45, 0.8, fill="white", radius=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_logo(slide, prs, 0.76, 0.52, 2.1)
    add_text(slide, "汇报完毕，感谢聆听", 2.55, 2.35, 8.3, 0.72, font_size=32, color="white", bold=True, align=PP_ALIGN.CENTER, font=FONT_SERIF)
    add_line(slide, 4.1, 3.35, 9.25, 3.35, color="gold", width=2.2)
    add_text(slide, "敬请各位老师批评指正", 3.7, 3.72, 5.9, 0.42, font_size=18, color="white", align=PP_ALIGN.CENTER)
    add_text(slide, "汇报人：沈秋雨 | 华中师范大学信息管理学院", 4.05, 5.75, 5.4, 0.34, font_size=12, color="white", align=PP_ALIGN.CENTER)
    return slide


def build() -> None:
    prs = Presentation(TEMPLATE)
    delete_all_slides(prs)

    title_slide(prs)
    contents_slide(prs)
    background_slide(prs)
    overall_idea_slide(prs)
    four_layers_slide(prs)
    decomposer_slide(prs)
    rationale_slide(prs)
    mapping_slide(prs)
    dataset_slide(prs)
    results_slide(prs)
    conclusion_slide(prs)
    ending_slide(prs)

    prs.save(OUT)


if __name__ == "__main__":
    build()
    print(OUT)
